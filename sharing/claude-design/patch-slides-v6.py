"""Round 6 — in-place rewrite of Slide 27.

策略（安全版）：
  不動 slide 順序，只把原 #27 Conclusion 頁「清空所有 shape → 重建」
  成新的 thesis 頁。原 #28 Pick-the-tool 頁保持不動，因此 flow 變成：
    #26 Outcome
    #27 [REWRITTEN] thesis page (BEFORE WE PICK THE TOOL)
    #28 Pick the tool by the job (不動)
  原 #27 的 Rasmus 引言吸收進新 #27 底部。

頁數保持 33。避開 add_slide / delete_slide 的 rId duplicate bug。

Run in place:
  python3 patch-slides-v6.py
"""
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "claude-design-sharing.pptx")

TEMPLATE_BLUE = RGBColor(0x1F, 0x4E, 0x79)
TEMPLATE_GRAY = RGBColor(0x55, 0x66, 0x77)
TEMPLATE_DARK = RGBColor(0x15, 0x32, 0x41)
ACCENT_LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
MUTED_GRAY    = RGBColor(0x8A, 0x92, 0x9A)


def clear_all_shapes(slide):
    """Remove every shape on the slide (including title/body placeholders).
    Uses low-level XML removal to avoid python-pptx ShapeTree limitations."""
    spTree = slide.shapes._spTree
    # Collect all sp / pic / graphicFrame / grpSp / cxnSp children
    for child in list(spTree):
        tag = child.tag
        # Keep only the non-shape metadata (nvGrpSpPr, grpSpPr)
        if tag.endswith('}nvGrpSpPr') or tag.endswith('}grpSpPr'):
            continue
        spTree.remove(child)


def build_new_slide27(slide):
    """Thesis page: 'Low barrier · narrower range — pick by scenario' + limitation table + Rasmus quote."""

    # ===== Eyebrow =====
    eb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8), Inches(0.35))
    tf = eb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "BEFORE WE PICK THE TOOL"
    r.font.name = "Arial"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = TEMPLATE_BLUE

    # ===== Headline =====
    hl = slide.shapes.add_textbox(Inches(0.6), Inches(0.85), Inches(12.1), Inches(1.1))
    tf = hl.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Low barrier. Narrower range."
    r.font.name = "Arial"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = TEMPLATE_DARK

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r = p2.add_run()
    r.text = "Pick the tool by scenario, not by preference."
    r.font.name = "Arial"
    r.font.size = Pt(20)
    r.font.color.rgb = TEMPLATE_GRAY

    # ===== Body explanation (Chinese) =====
    body = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(12.1), Inches(0.75))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Claude Design 的 GUI 更易上手，但 Claude Code 依然涵蓋較多場景。"
    r.font.name = "Noto Sans TC"
    r.font.size = Pt(15)
    r.font.color.rgb = TEMPLATE_DARK

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(3)
    r = p2.add_run()
    r.text = "下方是結構性差異 — 這些是「依場景選工具」的依據。"
    r.font.name = "Noto Sans TC"
    r.font.size = Pt(13)
    r.font.color.rgb = TEMPLATE_GRAY

    # ===== Limitation table (結構性層級, 跟 #28 操作性層級錯開) =====
    rows = 6
    cols = 2
    tbl_left = Inches(0.6)
    tbl_top  = Inches(3.25)
    tbl_w    = Inches(12.1)
    tbl_h    = Inches(3.05)

    table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h)
    tbl = table_shape.table
    tbl.columns[0].width = Inches(4.8)
    tbl.columns[1].width = Inches(7.3)

    headers = ["Claude Design 做不到的", "影響的場景"]
    for c, text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEMPLATE_BLUE
        cell.margin_left = Inches(0.14)
        cell.margin_right = Inches(0.14)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)

    data = [
        ("Team shared skills / documents",
         "跨人、跨專案共用規則做不了 — 每人重貼"),
        ("Workspace-level instruction",
         "組織級 convention 無法持久注入"),
        ("Skill runner",
         "無法跑自動化檢查、lint、commit、sync"),
        ("MCP / hooks",
         "無法接 Jira / Teams / cron 自動化需求流"),
        ("Decision trace",
         "無法稽核、無 git blame、決策不可追"),
    ]

    for r_idx, (limit, impact) in enumerate(data, start=1):
        for c_idx, text in enumerate([limit, impact]):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = text
            if c_idx == 0:
                r.font.name = "Arial"
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = TEMPLATE_DARK
            else:
                r.font.name = "Noto Sans TC"
                r.font.size = Pt(11)
                r.font.color.rgb = TEMPLATE_DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_LIGHT if r_idx % 2 == 1 else WHITE
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)

    tbl.rows[0].height = Inches(0.42)
    for i in range(1, rows):
        tbl.rows[i].height = Inches((tbl_h.inches - 0.42) / (rows - 1))

    # ===== Rasmus quote at bottom (absorbed from old #27) =====
    qb = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.85))
    qtf = qb.text_frame
    qtf.word_wrap = True
    qtf.margin_left = 0

    qp = qtf.paragraphs[0]
    qp.alignment = PP_ALIGN.LEFT
    qr = qp.add_run()
    qr.text = "\u201cIt\u2019s not a design tool as much as it\u2019s a design production tool. "
    qr.font.name = "Arial"
    qr.font.size = Pt(12)
    qr.font.italic = True
    qr.font.color.rgb = TEMPLATE_GRAY

    qr2 = qp.add_run()
    qr2.text = "Design is methodical; like architecture, not like art.\u201d"
    qr2.font.name = "Arial"
    qr2.font.size = Pt(12)
    qr2.font.italic = True
    qr2.font.color.rgb = TEMPLATE_GRAY

    cp = qtf.add_paragraph()
    cp.alignment = PP_ALIGN.LEFT
    cp.space_before = Pt(3)
    cr = cp.add_run()
    cr.text = "— Rasmus Andersson (@rsms)"
    cr.font.name = "Arial"
    cr.font.size = Pt(10)
    cr.font.color.rgb = MUTED_GRAY


def main():
    prs = Presentation(SRC)
    print(f"Slide count: {len(prs.slides)}")

    # Old #27 is at index 26
    slide27 = prs.slides[26]

    # Step 1: Clear all shapes on Slide 27
    clear_all_shapes(slide27)
    print(f"  [OK] cleared all shapes on slide 27")

    # Step 2: Build new thesis content
    build_new_slide27(slide27)
    print(f"  [OK] rebuilt slide 27: thesis + limitation table + Rasmus quote")

    prs.save(SRC)
    print(f"\nSaved: {SRC}")


if __name__ == "__main__":
    main()
