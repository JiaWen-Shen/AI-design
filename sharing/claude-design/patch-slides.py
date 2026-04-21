"""Patch only slides 3, 29, 32 of the current claude-design-sharing.pptx.

Preserves user's manual edits to all other slides. Uses the shared helpers
from build-claude-design-sharing.py.

Run:
  cd .../sharing/claude-design/
  python3 patch-slides.py
"""
import os
import sys

# Reuse helpers from the main build script
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from importlib.machinery import SourceFileLoader
bm = SourceFileLoader(
    "buildmod",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "build-claude-design-sharing.py")
).load_module()

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "claude-design-sharing.pptx")
DST = SRC  # in-place

# palette shortcuts
RED = bm.RED; NAVY = bm.NAVY; GRAY = bm.GRAY
LGT = bm.LGT; WHITE = bm.WHITE; SOFT_RED = bm.SOFT_RED
emu = bm.emu

SLIDE_W = bm.SLIDE_W; SLIDE_H = bm.SLIDE_H

# ---------- Slide #3 · TOC as horizontal row list ----------
def build_toc_horizontal(slide):
    bm.clear_slide_shapes(slide)
    # Eyebrow + title
    bm.add_text(slide, "WHAT WE'LL COVER", emu(0.7), emu(0.5), emu(6), emu(0.4),
                size=14, bold=True, color=RED, letter_spacing=300)
    bm.add_text(slide, "Table of Contents", emu(0.7), emu(0.9), emu(10), emu(0.9),
                size=40, bold=True, color=NAVY)

    # 6 horizontal rows — each = one full-width row
    items = [
        ("Q1", "What's the difference between Claude Code and Claude Design?",
         "The GUI-vs-CLI story"),
        ("Q2", "How can Claude Design help us?",
         "Pain points it solves today"),
        ("Q3", "Does Claude Design improve delivery quality?",
         "Where it helps · where it doesn't"),
        ("Q4", "Do we still need Figma?",
         "Design tool vs. production tool"),
        ("※", "Hands-on walkthrough",
         "A designer's first run"),
        ("→", "What's next",
         "Long-term outcomes"),
    ]
    row_top = emu(2.2)
    row_h = emu(0.75)
    row_gap = emu(0.05)
    x = emu(0.7)
    row_w = emu(11.9)
    for i, (badge, title, sub) in enumerate(items):
        y = row_top + i * (row_h + row_gap)
        # row background (alternating for readability)
        fill = LGT if i % 2 == 0 else WHITE
        bm.add_rounded_rect(slide, x, y, row_w, row_h, fill=fill)
        # badge on left
        is_q = badge.startswith("Q")
        badge_fill = RED if is_q else NAVY
        bm.add_rounded_rect(slide, x + emu(0.2), y + emu(0.17), emu(0.7), emu(0.4), fill=badge_fill)
        bm.add_text(slide, badge, x + emu(0.2), y + emu(0.17), emu(0.7), emu(0.4),
                    size=14, bold=True, color=WHITE, align='center', valign='middle')
        # title
        bm.add_text(slide, title, x + emu(1.1), y + emu(0.08), emu(7.2), emu(0.65),
                    size=16, bold=True, color=NAVY, valign='middle')
        # sub (right side)
        bm.add_text(slide, sub, x + emu(8.4), y + emu(0.08), emu(3.3), emu(0.65),
                    size=13, color=GRAY, valign='middle')

    # footer
    bm.add_text(slide, "Claude Design — A Designer's Review",
                emu(0.7), emu(6.95), emu(8), emu(0.3),
                size=10, color=GRAY)

# ---------- Slide #29 · Conclusion 2 with Claude Design limitations ----------
def build_conclusion_2_with_limits(slide):
    bm.clear_slide_shapes(slide)
    bm.add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

    # Header
    bm.add_text(slide, "CONCLUSION · FOR OUR SCENARIO",
                emu(0.7), emu(0.35), emu(10), emu(0.4),
                size=14, bold=True, color=RED)
    bm.add_text(slide, "Pick the tool by the job.",
                emu(0.7), emu(0.85), emu(12), emu(0.75),
                size=32, bold=True, color=NAVY)

    # Two cards side-by-side
    card_top = emu(1.85)
    card_h = emu(3.75)
    card_w = emu(5.65)
    gap = emu(0.3)
    x1 = emu(0.7); x2 = x1 + card_w + gap

    # --- Card A: Claude Design main stage ---
    bm.add_rounded_rect(slide, x1, card_top, card_w, card_h, fill=LGT)
    bm.add_rounded_rect(slide, x1 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35), fill=NAVY)
    bm.add_text(slide, "A", x1 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35),
                size=14, bold=True, color=WHITE, align='center', valign='middle')
    bm.add_text(slide, "Pure visual validation · cross-team review",
                x1 + emu(1.05), card_top + emu(0.3), card_w - emu(1.4), emu(0.4),
                size=13, bold=True, color=GRAY)
    bm.add_text(slide, "Claude Design\nis the main stage.",
                x1 + emu(0.35), card_top + emu(0.9), card_w - emu(0.7), emu(1.2),
                size=24, bold=True, color=NAVY)
    bullets_a = [
        "Org-internal share URL for stakeholder review",
        "Inline comments · live canvas iteration",
        "One-click export to PDF / PPTX / Canva",
    ]
    for i, b in enumerate(bullets_a):
        by = card_top + emu(2.25 + i*0.35)
        bm.add_rounded_rect(slide, x1 + emu(0.4), by + emu(0.08), emu(0.12), emu(0.12), fill=NAVY)
        bm.add_text(slide, b, x1 + emu(0.65), by, card_w - emu(0.95), emu(0.3),
                    size=12, color=NAVY)

    # --- Card B: Claude Design's limitations (push you toward Claude Code) ---
    bm.add_rounded_rect(slide, x2, card_top, card_w, card_h, fill=SOFT_RED)
    bm.add_rounded_rect(slide, x2 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35), fill=RED)
    bm.add_text(slide, "B", x2 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35),
                size=14, bold=True, color=WHITE, align='center', valign='middle')
    bm.add_text(slide, "Ship to RD · automate · version-control",
                x2 + emu(1.05), card_top + emu(0.3), card_w - emu(1.4), emu(0.4),
                size=13, bold=True, color=RED)
    bm.add_text(slide, "Claude Design's limits\npush you to Claude Code.",
                x2 + emu(0.35), card_top + emu(0.9), card_w - emu(0.7), emu(1.2),
                size=22, bold=True, color=NAVY)
    # 5 limitation bullets — mined from claude-design-vs-cli-designer-handoff.md
    limits = [
        ("Design system", "No auto-sync — updates need manual Remix"),
        ("Version control", "No rollback · tokens are a black box · can't diff"),
        ("Requirements intake", "Every update pasted by hand (no MCP / cron)"),
        ("Delivery to RD", "Two translations (visual → code → prod)"),
        ("Figma export", "No official path — third-party tools required"),
    ]
    for i, (k, v) in enumerate(limits):
        by = card_top + emu(2.2 + i*0.3)
        bm.add_rounded_rect(slide, x2 + emu(0.4), by + emu(0.08), emu(0.12), emu(0.12), fill=RED)
        bm.add_text(slide, k, x2 + emu(0.65), by, emu(1.6), emu(0.28),
                    size=11, bold=True, color=RED)
        bm.add_text(slide, v, x2 + emu(2.25), by, card_w - emu(2.55), emu(0.28),
                    size=11, color=NAVY)

    # Relay callout strip at bottom
    relay_y = emu(5.85)
    relay_h = emu(1.1)
    bm.add_rounded_rect(slide, emu(0.7), relay_y, emu(11.9), relay_h, fill=NAVY)
    bm.add_text(slide, "It's not either / or — it's a relay.",
                emu(1.0), relay_y + emu(0.18), emu(5.8), emu(0.45),
                size=20, bold=True, color=WHITE)
    bm.add_text(slide, "Two tools. One chain. Different jobs.",
                emu(1.0), relay_y + emu(0.65), emu(5.8), emu(0.35),
                size=13, color=RGBColor(0xB8, 0xC3, 0xC9))
    # chain pills
    chain_y = relay_y + emu(0.38)
    pill_w = emu(1.3); pill_h = emu(0.38); pill_gap = emu(0.35)
    cx1 = emu(7.3)
    cx2 = cx1 + pill_w + pill_gap
    cx3 = cx2 + pill_w + pill_gap
    for label, x, color in [
        ("Design", cx1, RED),
        ("Code", cx2, RGBColor(0x30, 0x5A, 0x72)),
        ("RD", cx3, GRAY),
    ]:
        bm.add_rounded_rect(slide, x, chain_y, pill_w, pill_h, fill=color)
        bm.add_text(slide, label, x, chain_y, pill_w, pill_h,
                    size=13, bold=True, color=WHITE, align='center', valign='middle')
    for ax in [cx1 + pill_w, cx2 + pill_w]:
        bm.add_text(slide, "→", ax, chain_y, pill_gap, pill_h,
                    size=18, color=WHITE, align='center', valign='middle')

# ---------- Slide #32 · References as plain numbered list ----------
def build_references_plain(slide):
    bm.clear_slide_shapes(slide)
    bm.add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    # Header
    bm.add_text(slide, "APPENDIX", emu(0.7), emu(0.5), emu(6), emu(0.4),
                size=13, bold=True, color=RED, letter_spacing=400)
    bm.add_text(slide, "References", emu(0.7), emu(0.9), emu(10), emu(0.8),
                size=36, bold=True, color=NAVY)

    # Plain markdown-like numbered list, one per row
    refs = [
        "Introducing Claude Design by Anthropic Labs — anthropic.com/news/claude-design-anthropic-labs",
        "Claude (@claudeai) launch post on X — x.com/claudeai/status/2045156267690213649",
        "Ryan Mather (@Flomerboy) — Tips for Claude Design — x.com/Flomerboy/status/2045162321589252458",
        "Claude Design — claude.ai/design",
        "Figma Skills: Code to design workflows — figma.com/community/skills",
        "Rasmus Andersson (@rsms) — Quote — x.com/rsms/status/2045239193971179851",
    ]
    row_top = emu(2.0)
    row_h = emu(0.58)
    for i, text in enumerate(refs):
        y = row_top + i * row_h
        bm.add_text(slide, f"{i+1}.  {text}",
                    emu(0.7), y, emu(12.0), row_h,
                    size=15, color=NAVY, valign='middle')

# ---------- Main ----------
def main():
    prs = Presentation(SRC)
    total = len(prs.slides)
    print(f"Loaded: {total} slides")
    assert total >= 32, f"expected >=32 slides, got {total}"

    # Slide 3 (index 2) — TOC
    build_toc_horizontal(prs.slides[2])
    print("  patched slide 3 (TOC → horizontal rows)")

    # Slide 29 (index 28) — Conclusion 2 with limitations
    build_conclusion_2_with_limits(prs.slides[28])
    print("  patched slide 29 (Conclusion 2 → with Claude Design limitations)")

    # Slide 32 (index 31) — References plain list
    build_references_plain(prs.slides[31])
    print("  patched slide 32 (References → plain numbered list)")

    prs.save(DST)
    print(f"Saved: {DST}")

if __name__ == "__main__":
    main()
