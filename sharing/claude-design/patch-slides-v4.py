"""Round 4 — final pre-share polish pass.

Applies all the fixes the user accepted from the editorial review:

  #1  "By Anthropic Labs" → "From Anthropic Labs"
  #2  "Prototype to Production" → "From Idea to Artifact"
  #3  "Let‘s start" → "Let's start" (straight apostrophe)
  #8  "Keep the output quality" → "Same engine, lower barrier."
  #11 Row 2 left: "AI hard to harness" → "Zero-to-one scaffolding"
  #11 Row 2 right body: updated to "Design system auto-built from codebase + design files."
  #13 caption "v.s" → "Low barrier · High ceiling" replacement
  #13 WHY closing comma tightening
  #15 EXPORT label honest about Claude Code as the bridge
  #19 "Upload codes" → "Upload code" + strip non-breaking space
  #27 three benefit bullets reworded to what ships today
  #28 eyebrow space + "FALL BACK" → "REACH FOR" + "real-time co-edit" → safer
  #29 "Cross-person & team handover" → "Cross-team handover"
  #30 strip vertical-tab (\\x0b) inside Ryan quote

Run in place:
  python3 patch-slides-v4.py
"""
import os
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "claude-design-sharing.pptx")

# Text-level find/replace operations, applied across every run on every slide.
# Each entry: (slide_index_or_None, old, new). If slide_index is None, applies
# globally (first match on every matching slide).
REPLACEMENTS = [
    # --- #1 cover ---
    (0, "By Anthropic Labs.", "From Anthropic Labs."),
    # --- #2 briefing ---
    (1, "Prototype to Production", "From Idea to Artifact"),
    # --- #3 TOC subtitle (curly-left → straight) ---
    (2, "Let\u2018s start", "Let's start"),
    # --- #8 A2 callout subtitle ---
    (7, "Keep the output quality.", "Same engine, lower barrier."),
    # --- #11 "What Claude Design Solved" row 2 ---
    (10, "AI hard to harness", "Zero-to-one scaffolding"),
    (10, "Pre-built design system\nFrom codebase + design files.",
         "Design system auto-built\nfrom codebase + design files."),
    # --- #13 caption + WHY closing ---
    (12, "Easier to drive.     v.s      Wider range", "Low barrier  ·  High ceiling"),
    (12, "across time, or across teams.", "across time or across teams."),
    # --- #15 EXPORT arrow label ---
    (14, "EXPORT  \u2192  Claude Design output into Figma for final polish",
         "EXPORT  \u2192  Claude Code output  \u2192  Figma (for polish)"),
    # --- #19 "Upload codes" + non-breaking space ---
    (18, "Upload codes", "Upload code"),
    (18, "Figma\xa0", "Figma"),
    # --- #27 three benefit bullets ---
    (26, "Turn ideas into video demos. Pull from docs and meetings to iterate on design. Build one-off tools on demand.",
         "Quick motion mockups via screen-capture. Paste meeting notes · docs → fast design iterations. One-off visual artifacts for decks / pitches."),
    # --- #28 Conclusion 2 ---
    (27, "Better harness\u00b7 version-control", "Better harness  \u00b7  version-control"),
    (27, "B · FALL BACK TO CLAUDE CODE WHEN", "B · REACH FOR CLAUDE CODE WHEN"),
    (27, "Team-project workspaces — shared folders, real-time co-edit (Team / Enterprise plan)",
         "Team-project workspaces — shared folders · inline comments (Team / Enterprise plan)"),
    # --- #29 What's next ---
    (28, "Cross-person & team handover", "Cross-team handover"),
    # --- #30 Ryan quote vertical tab ---
    (29, "\x0b", "  "),
]


def apply_run_level(slide, old, new):
    """Walk every text run on the slide. If run.text contains old, replace it.
    Preserves run formatting."""
    replaced = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            # Collect the paragraph text joined from its runs so we can detect
            # matches spanning runs. If the match is entirely inside one run,
            # just replace that run's text.
            full = "".join(r.text for r in para.runs)
            if old not in full:
                continue
            # Try single-run replacement first (preserves formatting best).
            single_run_hit = False
            for r in para.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    replaced += 1
                    single_run_hit = True
                    break
            if single_run_hit:
                continue
            # Match spans runs — put the whole replacement into the first run,
            # blank the rest in the matched span. This is rare for our cases.
            first_run = para.runs[0]
            new_full = full.replace(old, new, 1)
            first_run.text = new_full
            for r in para.runs[1:]:
                r.text = ""
            replaced += 1
    return replaced


def main():
    prs = Presentation(SRC)
    for idx, old, new in REPLACEMENTS:
        slide = prs.slides[idx]
        n = apply_run_level(slide, old, new)
        status = "OK" if n > 0 else "MISS"
        print(f"  [{status}] slide {idx+1}: {old[:60]!r} → {new[:60]!r}  ({n} hit)")
    prs.save(SRC)
    print(f"\nSaved: {SRC}")


if __name__ == "__main__":
    main()
