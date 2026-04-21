"""Rebuild claude-design-sharing.pptx per draft 2 (31 slides).

Reads claude-design-sharing.v1.pptx, inserts 11 new slides with native pptx
shapes (so the user can edit icons/text directly in PowerPoint), drops old #3
and #18 (their content is subsumed by the new Q&A flow), and reorders to match
draft 2.

Run:
  cd .../sharing/claude-design/
  python3 build-claude-design-sharing.py
"""
import os
from copy import deepcopy
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "claude-design-sharing.v1.pptx")
DST = os.path.join(HERE, "claude-design-sharing.pptx")

# ---------- TrendLife palette ----------
RED   = RGBColor(0xD7, 0x19, 0x20)
NAVY  = RGBColor(0x15, 0x32, 0x41)
GRAY  = RGBColor(0x78, 0x80, 0x86)
LGT   = RGBColor(0xF4, 0xF6, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_RED = RGBColor(0xFB, 0xE5, 0xE6)
DARK_TERM = RGBColor(0x0F, 0x17, 0x24)
TERM_GREEN = RGBColor(0x8F, 0xE8, 0xA8)
TERM_DIM = RGBColor(0x5B, 0x6B, 0x85)

# 16:9, 13.33 x 7.5 inches = 12192000 x 6858000 EMU
SLIDE_W = 12192000
SLIDE_H = 6858000

def emu(inches_float):
    return int(inches_float * 914400)

# ---------- Shape helpers ----------
def _set_solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def _set_no_line(shape):
    shape.line.fill.background()

def _set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)

def add_rect(slide, left, top, width, height, fill=None, line=None, line_pt=1.0, shape_type=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        _set_solid_fill(s, fill)
    else:
        s.fill.background()
    if line is None:
        _set_no_line(s)
    else:
        _set_line(s, line, line_pt)
    return s

def add_rounded_rect(slide, left, top, width, height, fill=None, line=None, line_pt=1.0):
    return add_rect(slide, left, top, width, height, fill=fill, line=line, line_pt=line_pt,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

def add_oval(slide, cx, cy, r, fill=None, line=None, line_pt=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    if fill is not None:
        _set_solid_fill(s, fill)
    else:
        s.fill.background()
    if line is None:
        _set_no_line(s)
    else:
        _set_line(s, line, line_pt)
    return s

def add_line(slide, x1, y1, x2, y2, rgb=NAVY, width_pt=1.0, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = rgb
    c.line.width = Pt(width_pt)
    if dash:
        # 'dash', 'round_dot' etc via XML
        lnEl = c.line._get_or_add_ln()
        prstDash = etree.SubElement(lnEl, qn('a:prstDash'))
        prstDash.set('val', dash)
    return c

def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=NAVY, align='left', valign='top',
             font='Inter', letter_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    if valign == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    # First paragraph already exists
    p = tf.paragraphs[0]
    if align == 'center': p.alignment = PP_ALIGN.CENTER
    elif align == 'right': p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Explicitly disable underline — prevents LibreOffice from rendering
    # inherited underline from master/layout run properties.
    run.font.underline = False
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(letter_spacing))
    return tb

# ---------- Slide copy / reorder utilities ----------
def add_blank_slide(prs, layout_idx=7):
    """Default to 'Black Blank' layout (index 7). Disable layout/master shape
    inheritance so our new slides render cleanly without template graphics
    bleeding through (the layout has a decorative rectangle + graphic that
    cause a red-line artifact in LibreOffice around y=4.7in)."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Suppress inheritance of layout/master non-placeholder shapes:
    # <p:sld showMasterSp="0"> at the root of the slide XML.
    slide._element.set('showMasterSp', '0')
    return slide

def clear_slide_shapes(slide):
    """Remove all shapes from a slide (keeping only placeholders from layout)."""
    sp_tree = slide.shapes._spTree
    for shp in list(slide.shapes):
        sp_tree.remove(shp._element)

def reorder_slides(prs, new_order):
    """Reorder slides. new_order = list of indices into the current order."""
    sldIdLst = prs.slides._sldIdLst
    sldId_list = list(sldIdLst)
    for el in sldId_list:
        sldIdLst.remove(el)
    for idx in new_order:
        sldIdLst.append(sldId_list[idx])

# ---------- Slide builders ----------
def build_toc(slide):
    clear_slide_shapes(slide)
    # eyebrow
    add_text(slide, "WHAT WE'LL COVER", emu(0.7), emu(0.5), emu(6), emu(0.4),
             size=14, bold=True, color=RED, letter_spacing=300)
    # title
    add_text(slide, "Table of Contents", emu(0.7), emu(0.9), emu(8), emu(0.9),
             size=40, bold=True, color=NAVY)
    # 4 Q cards across the top
    qs = [
        ("Q1", "What's different from\nClaude Code?", "The GUI-vs-CLI story"),
        ("Q2", "How does Claude Design\nhelp us?", "Pain points it solves today"),
        ("Q3", "Does it improve\ndelivery quality?", "Where it helps, where it doesn't"),
        ("Q4", "Do we still need\nFigma?", "Design tool vs. production tool"),
    ]
    top = emu(2.4)
    card_w = emu(2.85)
    card_h = emu(2.4)
    gap = emu(0.15)
    start_x = emu(0.7)
    for i, (q, title, sub) in enumerate(qs):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(slide, x, top, card_w, card_h, fill=LGT)
        # Q badge
        add_rounded_rect(slide, x + emu(0.2), top + emu(0.2), emu(0.75), emu(0.4), fill=RED)
        add_text(slide, q, x + emu(0.2), top + emu(0.2), emu(0.75), emu(0.4),
                 size=14, bold=True, color=WHITE, align='center', valign='middle')
        # title
        add_text(slide, title, x + emu(0.2), top + emu(0.8), card_w - emu(0.4), emu(1.0),
                 size=17, bold=True, color=NAVY)
        # sub
        add_text(slide, sub, x + emu(0.2), top + emu(1.75), card_w - emu(0.4), emu(0.6),
                 size=12, color=GRAY)
    # Two pills underneath
    pills = [("Hands-on walkthrough", "A designer's first run"),
             ("What's next", "Long-term outcomes")]
    pill_top = emu(5.2)
    pill_w = emu(5.85)
    pill_h = emu(1.1)
    pill_gap = emu(0.3)
    for i, (title, sub) in enumerate(pills):
        x = start_x + i * (pill_w + pill_gap)
        add_rounded_rect(slide, x, pill_top, pill_w, pill_h, fill=SOFT_RED)
        add_text(slide, title, x + emu(0.3), pill_top + emu(0.15), pill_w - emu(0.6), emu(0.5),
                 size=18, bold=True, color=NAVY)
        add_text(slide, sub, x + emu(0.3), pill_top + emu(0.65), pill_w - emu(0.6), emu(0.4),
                 size=13, color=GRAY)
    # footer
    add_text(slide, "Claude Design — A Designer's Review", emu(0.7), emu(6.8), emu(8), emu(0.35),
             size=10, color=GRAY)

def build_subtitle(slide):
    clear_slide_shapes(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_text(slide, "SECTION 01", emu(0.9), emu(2.5), emu(6), emu(0.4),
             size=14, bold=True, color=RED, letter_spacing=400)
    add_text(slide, "Let's start with the questions\nmost people care about.",
             emu(0.9), emu(3.1), emu(11), emu(2.5),
             size=48, bold=True, color=WHITE)
    add_text(slide, "Four questions. Four honest answers.",
             emu(0.9), emu(5.6), emu(11), emu(0.5),
             size=20, color=RGBColor(0xB8, 0xC3, 0xC9))

def build_q_prompt(slide, q_num, question_text, eyebrow_label=None):
    clear_slide_shapes(slide)
    # light bg
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    # Big oversized Q numeral on the left (split letter + digit so no wrap)
    add_text(slide, q_num[0], emu(0.3), emu(1.4), emu(3.2), emu(5),
             size=260, bold=True, color=SOFT_RED, font='Inter',
             align='center')
    add_text(slide, q_num[1:], emu(2.5), emu(1.4), emu(2.8), emu(5),
             size=260, bold=True, color=SOFT_RED, font='Inter',
             align='center')
    # Eyebrow
    eyebrow = eyebrow_label or f"QUESTION {q_num[1:].zfill(2)}"
    add_text(slide, eyebrow, emu(5.5), emu(2.6), emu(7), emu(0.4),
             size=14, bold=True, color=RED, letter_spacing=400)
    # Title
    add_text(slide, question_text, emu(5.5), emu(3.1), emu(7.3), emu(3),
             size=34, bold=True, color=NAVY)

def build_a1_callout(slide):
    """A1: Claude Design : Claude Code = GitHub : Git (native shapes).
    Top row = analogy (Git / GitHub), bottom row = tools (Claude Code / Claude Design).
    """
    clear_slide_shapes(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

    # Eyebrow + headline
    add_text(slide, "ANSWER 01 · IN ONE LINE", emu(0.7), emu(0.4), emu(10), emu(0.4),
             size=13, bold=True, color=RED, letter_spacing=400)
    add_text(slide, "Claude Design : Claude Code  =  GitHub : Git",
             emu(0.7), emu(0.85), emu(12), emu(0.9),
             size=34, bold=True, color=NAVY)
    add_text(slide, "A GUI layer that directs the same underlying work.",
             emu(0.7), emu(1.85), emu(12), emu(0.5),
             size=18, color=GRAY)

    # Two-row comparison
    #  Top row: Git (CLI)  |  GitHub (GUI)
    #  Bottom row: Claude Code (CLI)  |  Claude Design (GUI)
    row_h = emu(2.1)
    col_w = emu(5.6)
    gap = emu(0.5)
    start_x = emu(0.7)
    top1 = emu(2.7)
    top2 = emu(4.95)

    def panel_label(x, y, label_text, color=NAVY):
        # letter_spacing removed — LibreOffice renders spc attr with a red underline artifact
        add_text(slide, label_text, x, y, col_w, emu(0.35),
                 size=13, bold=True, color=color)

    def terminal_panel(x, y, commands):
        # dark terminal
        t = add_rounded_rect(slide, x, y, col_w, row_h, fill=DARK_TERM)
        # 3 dots
        for i, c in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
            add_oval(slide, x + emu(0.25 + i*0.22), y + emu(0.22), emu(0.08), fill=c)
        # commands
        for i, cmd in enumerate(commands):
            col = TERM_GREEN if i < 2 else TERM_DIM
            add_text(slide, cmd, x + emu(0.25), y + emu(0.55 + i*0.32), col_w - emu(0.4), emu(0.3),
                     size=13, color=col, font='Menlo')

    def gui_panel(x, y, kind):
        # light card
        t = add_rounded_rect(slide, x, y, col_w, row_h, fill=WHITE, line=NAVY, line_pt=2)
        # title bar
        add_rounded_rect(slide, x, y, col_w, emu(0.35), fill=LGT)
        for i, c in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
            add_oval(slide, x + emu(0.2 + i*0.18), y + emu(0.18), emu(0.07), fill=c)
        if kind == 'github':
            # commit list with dots + labels
            add_oval(slide, x + emu(0.45), y + emu(0.75), emu(0.11), fill=RED)
            add_text(slide, "fix layout bug", x + emu(0.7), y + emu(0.58), emu(3), emu(0.3), size=12, bold=True, color=NAVY)
            add_text(slide, "Karen · 2 min ago", x + emu(0.7), y + emu(0.85), emu(3), emu(0.25), size=10, color=GRAY)
            add_oval(slide, x + emu(0.45), y + emu(1.2), emu(0.11), fill=NAVY)
            add_text(slide, "add header", x + emu(0.7), y + emu(1.03), emu(3), emu(0.3), size=12, bold=True, color=NAVY)
            add_text(slide, "Karen · 1 hr ago", x + emu(0.7), y + emu(1.3), emu(3), emu(0.25), size=10, color=GRAY)
            # Push button
            btn_w = emu(1.35); btn_h = emu(0.38)
            btnx = x + col_w - btn_w - emu(0.25)
            btny = y + row_h - btn_h - emu(0.25)
            add_rounded_rect(slide, btnx, btny, btn_w, btn_h, fill=RED)
            add_text(slide, "Push to origin", btnx, btny, btn_w, btn_h,
                     size=11, bold=True, color=WHITE, align='center', valign='middle')
        elif kind == 'claude_design':
            # chat sidebar (left 1.4")
            sidebar_w = emu(1.55)
            add_rect(slide, x, y + emu(0.35), sidebar_w, row_h - emu(0.35), fill=LGT)
            add_rounded_rect(slide, x + emu(0.2), y + emu(0.55), sidebar_w - emu(0.4), emu(0.3), fill=WHITE, line=GRAY, line_pt=0.75)
            add_rounded_rect(slide, x + emu(0.2), y + emu(0.9), sidebar_w - emu(0.4), emu(0.4), fill=RED)
            add_text(slide, "design pricing page", x + emu(0.2), y + emu(0.9), sidebar_w - emu(0.4), emu(0.4),
                     size=9, color=WHITE, align='center', valign='middle')
            add_rounded_rect(slide, x + emu(0.2), y + emu(1.35), sidebar_w - emu(0.4), emu(0.3), fill=WHITE, line=GRAY, line_pt=0.75)
            # canvas cards
            cx = x + sidebar_w + emu(0.2); cy = y + emu(0.55)
            cw = col_w - sidebar_w - emu(0.4)
            add_rounded_rect(slide, cx, cy, cw, emu(0.28), fill=LGT)
            card_w = (cw - emu(0.2)) / 3
            for i, fc in enumerate([SOFT_RED, LGT, LGT]):
                add_rounded_rect(slide, cx + i*(card_w + emu(0.1)), cy + emu(0.4), card_w, emu(1.0), fill=fc, line=NAVY if i==1 else None, line_pt=1.5)

    # Draw panels FIRST, labels LAST (z-order: labels on top).
    # Row 1: Git (CLI) → GitHub (GUI)
    terminal_panel(start_x, top1, ["$ git add .", "$ git commit -m \"fix\"", "$ git push origin main", "1 file changed"])
    gui_panel(start_x + col_w + gap, top1, 'github')
    # Row 2: Claude Code (CLI) → Claude Design (GUI)
    terminal_panel(start_x, top2, ["$ claude \"design a pricing page\"", "● Reading /App.jsx", "● Writing /Pricing.jsx", "✓ preview at localhost:3000"])
    gui_panel(start_x + col_w + gap, top2, 'claude_design')

    # Labels (drawn after panels so they're on top)
    panel_label(start_x, top1 - emu(0.32), "GIT   ·   CLI", color=GRAY)
    panel_label(start_x + col_w + gap, top1 - emu(0.32), "GITHUB   ·   GUI", color=GRAY)
    panel_label(start_x, top2 - emu(0.32), "CLAUDE CODE   ·   CLI", color=RED)
    panel_label(start_x + col_w + gap, top2 - emu(0.32), "CLAUDE DESIGN   ·   GUI", color=RED)

    # Arrows between panels
    mid_y = top1 + row_h/2
    add_text(slide, "→", start_x + col_w, mid_y - emu(0.25), gap, emu(0.5),
             size=28, color=GRAY, align='center', valign='middle')
    mid_y2 = top2 + row_h/2
    add_text(slide, "→", start_x + col_w, mid_y2 - emu(0.25), gap, emu(0.5),
             size=28, color=RED, align='center', valign='middle')

    # Left edge equivalence sign between rows
    add_text(slide, "≈", emu(6.2), top1 + row_h + emu(0.05), emu(1), emu(0.4),
             size=32, color=GRAY, align='center', valign='middle')

def build_a2_text_callout(slide):
    """A2: text-only callout (no image)."""
    clear_slide_shapes(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    add_text(slide, "ANSWER 02 · IN ONE LINE", emu(1), emu(1.3), emu(10), emu(0.4),
             size=14, bold=True, color=RED, letter_spacing=400)
    add_text(slide, "A designer-friendly alternative\nto Claude Code.",
             emu(1), emu(1.9), emu(11), emu(2.5),
             size=52, bold=True, color=NAVY)
    add_text(slide, "Lower the learning curve. Keep the output quality. Let designers ship prototypes without becoming terminal users.",
             emu(1), emu(4.7), emu(11), emu(1),
             size=20, color=GRAY)
    # 3 bullets
    bullets = [
        "No CLI. No Git. Just a canvas and a chat box.",
        "Pre-built scaffolding: design system import, design-decision prompts, export paths.",
        "The next 3 slides break down exactly which pain points it removes.",
    ]
    for i, b in enumerate(bullets):
        add_rounded_rect(slide, emu(1), emu(5.8 + i*0.45), emu(0.15), emu(0.15), fill=RED)
        add_text(slide, b, emu(1.3), emu(5.75 + i*0.45), emu(11), emu(0.4),
                 size=14, color=NAVY)

def build_a3_callout(slide):
    """A3: camera / gear / Claude analogy (3 rows, native shapes)."""
    clear_slide_shapes(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

    # Header
    add_text(slide, "ANSWER 03 · HONESTLY", emu(0.7), emu(0.3), emu(10), emu(0.35),
             size=13, bold=True, color=RED, letter_spacing=400)
    add_text(slide, "Yes — for structure.  No — for ceiling.",
             emu(0.7), emu(0.7), emu(12), emu(0.7),
             size=28, bold=True, color=NAVY)

    # 3-row analogy grid
    row_h = emu(1.55)
    left_cx = emu(3.3)    # center x of left analogy icon
    right_cx = emu(9.9)
    row_top1 = emu(1.9)   # cameras
    row_top2 = row_top1 + row_h + emu(0.1)
    row_top3 = row_top2 + row_h + emu(0.1)

    def analogy_row(y, left_builder, left_top, left_sub, right_builder, right_top, right_sub):
        # left background
        add_rounded_rect(slide, left_cx - emu(1.5), y, emu(3.0), row_h, fill=LGT)
        left_builder(left_cx, y + emu(0.1))
        add_text(slide, left_top, left_cx - emu(1.5), y + row_h - emu(0.55), emu(3.0), emu(0.3),
                 size=13, bold=True, color=NAVY, align='center')
        add_text(slide, left_sub, left_cx - emu(1.5), y + row_h - emu(0.28), emu(3.0), emu(0.25),
                 size=10, color=GRAY, align='center')
        # vertical separator
        add_line(slide, emu(6.6), y + emu(0.2), emu(6.6), y + row_h - emu(0.2), rgb=GRAY, width_pt=1, dash='dash')
        # right background
        add_rounded_rect(slide, right_cx - emu(1.5), y, emu(3.0), row_h, fill=LGT)
        right_builder(right_cx, y + emu(0.1))
        add_text(slide, right_top, right_cx - emu(1.5), y + row_h - emu(0.55), emu(3.0), emu(0.3),
                 size=13, bold=True, color=NAVY, align='center')
        add_text(slide, right_sub, right_cx - emu(1.5), y + row_h - emu(0.28), emu(3.0), emu(0.25),
                 size=10, color=GRAY, align='center')

    # --- Row 1 icon builders ---
    def point_and_shoot(cx, cy):
        body = add_rect(slide, cx - emu(0.8), cy, emu(1.6), emu(0.9), fill=NAVY)
        add_oval(slide, cx, cy + emu(0.45), emu(0.32), fill=LGT)
        add_oval(slide, cx, cy + emu(0.45), emu(0.2), fill=NAVY)
        add_oval(slide, cx, cy + emu(0.45), emu(0.1), fill=GRAY)
        add_rect(slide, cx + emu(0.55), cy + emu(0.08), emu(0.12), emu(0.07), fill=RED)

    def dslr(cx, cy):
        # body
        add_rect(slide, cx - emu(1.0), cy, emu(2.0), emu(0.95), fill=NAVY)
        # grip
        add_rounded_rect(slide, cx - emu(1.0), cy + emu(0.2), emu(0.35), emu(0.75), fill=RGBColor(0x0A, 0x1A, 0x23))
        # viewfinder bump
        add_rect(slide, cx - emu(0.25), cy - emu(0.12), emu(0.4), emu(0.12), fill=NAVY)
        # lens (concentric)
        add_oval(slide, cx + emu(0.2), cy + emu(0.5), emu(0.5), fill=RGBColor(0x0A, 0x1A, 0x23))
        add_oval(slide, cx + emu(0.2), cy + emu(0.5), emu(0.4), fill=NAVY, line=GRAY, line_pt=0.5)
        add_oval(slide, cx + emu(0.2), cy + emu(0.5), emu(0.25), fill=RGBColor(0x0A, 0x1A, 0x23))
        add_oval(slide, cx + emu(0.2), cy + emu(0.5), emu(0.12), fill=GRAY)
        # dials
        add_oval(slide, cx - emu(0.6), cy + emu(0.13), emu(0.1), fill=GRAY)
        add_oval(slide, cx + emu(0.75), cy + emu(0.13), emu(0.1), fill=GRAY)
        add_rect(slide, cx + emu(0.55), cy + emu(0.28), emu(0.35), emu(0.06), fill=RED)

    # --- Row 2 icon builders ---
    def auto_gear(cx, cy):
        # label column P R N D
        for i, letter in enumerate(["P", "R", "N"]):
            add_text(slide, letter, cx - emu(0.9), cy + emu(0.08 + i*0.22), emu(0.25), emu(0.22),
                     size=12, bold=True, color=NAVY, align='center')
        add_text(slide, "D", cx - emu(0.9), cy + emu(0.74), emu(0.25), emu(0.22),
                 size=12, bold=True, color=RED, align='center')
        # stick body
        add_rounded_rect(slide, cx - emu(0.2), cy, emu(0.4), emu(1.0), fill=NAVY)
        # knob
        add_oval(slide, cx, cy + emu(0.3), emu(0.13), fill=RED)

    def manual_gear(cx, cy):
        # H-pattern at y+0.4~0.8
        top_y = cy + emu(0.35)
        bot_y = cy + emu(0.8)
        left_x = cx - emu(0.75)
        mid_x = cx
        right_x = cx + emu(0.75)
        # horizontal bars
        add_line(slide, left_x, top_y, right_x, top_y, rgb=NAVY, width_pt=3)
        add_line(slide, left_x, bot_y, right_x, bot_y, rgb=NAVY, width_pt=3)
        # vertical bar
        add_line(slide, mid_x, top_y, mid_x, bot_y, rgb=NAVY, width_pt=3)
        # dots + labels
        positions_top = [(left_x, "1"), (mid_x, "3"), (right_x, "5")]
        positions_bot = [(left_x, "2"), (mid_x, "4"), (right_x, "6")]
        for x, lab in positions_top:
            add_oval(slide, x, top_y, emu(0.1), fill=NAVY)
            add_text(slide, lab, x - emu(0.15), top_y - emu(0.35), emu(0.3), emu(0.22),
                     size=10, bold=True, color=NAVY, align='center')
        for x, lab in positions_bot:
            fill = RED if lab == "6" else NAVY
            add_oval(slide, x, bot_y, emu(0.1), fill=fill)
            add_text(slide, lab, x - emu(0.15), bot_y + emu(0.06), emu(0.3), emu(0.22),
                     size=10, bold=True, color=fill, align='center')

    # --- Row 3 icon builders: Claude Code vs Claude Design (shrunk to 0.8" height) ---
    def claude_code_icon(cx, cy):
        # terminal window — smaller to avoid overlapping labels
        add_rounded_rect(slide, cx - emu(0.9), cy, emu(1.8), emu(0.8), fill=DARK_TERM)
        for i, c in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
            add_oval(slide, cx - emu(0.78) + i*emu(0.14), cy + emu(0.1), emu(0.05), fill=c)
        add_text(slide, "$ claude design", cx - emu(0.82), cy + emu(0.22), emu(1.6), emu(0.22),
                 size=10, color=TERM_GREEN, font='Menlo')
        add_text(slide, "● designing…", cx - emu(0.82), cy + emu(0.42), emu(1.6), emu(0.22),
                 size=8, color=TERM_DIM, font='Menlo')
        add_text(slide, "● step 2 / 5", cx - emu(0.82), cy + emu(0.6), emu(1.6), emu(0.22),
                 size=8, color=TERM_DIM, font='Menlo')

    def claude_design_icon(cx, cy):
        # canvas window — smaller
        add_rounded_rect(slide, cx - emu(0.9), cy, emu(1.8), emu(0.8), fill=WHITE, line=NAVY, line_pt=1.5)
        add_rect(slide, cx - emu(0.9), cy + emu(0.08), emu(0.5), emu(0.72), fill=LGT)
        add_rounded_rect(slide, cx - emu(0.83), cy + emu(0.18), emu(0.35), emu(0.1), fill=WHITE, line=GRAY, line_pt=0.5)
        add_rounded_rect(slide, cx - emu(0.83), cy + emu(0.33), emu(0.35), emu(0.16), fill=RED)
        add_rounded_rect(slide, cx - emu(0.83), cy + emu(0.54), emu(0.35), emu(0.1), fill=WHITE, line=GRAY, line_pt=0.5)
        # canvas cards
        cx2 = cx - emu(0.3); cy2 = cy + emu(0.18)
        for i, fc in enumerate([SOFT_RED, LGT, LGT]):
            add_rounded_rect(slide, cx2 + i*emu(0.37), cy2, emu(0.28), emu(0.55), fill=fc, line=NAVY if i==1 else None, line_pt=1)

    analogy_row(row_top1, point_and_shoot, "Point-and-shoot", "傻瓜相機",
                dslr, "DSLR", "單眼")
    analogy_row(row_top2, auto_gear, "Automatic", "自排車",
                manual_gear, "Manual", "手排車")
    analogy_row(row_top3, claude_design_icon, "Claude Design", "GUI · easier to drive",
                claude_code_icon, "Claude Code", "CLI · wider range")

    # Caption strip at bottom
    cap_y = emu(6.8)
    add_rounded_rect(slide, emu(1.5), cap_y, emu(10.3), emu(0.5), fill=SOFT_RED)
    add_text(slide, "Easier to drive.   ·   Wider range.",
             emu(1.5), cap_y, emu(10.3), emu(0.5),
             size=16, bold=True, color=NAVY, align='center', valign='middle')

def build_a4_callout(slide):
    """A4: Figma MCP bidirectional flow (native shapes)."""
    clear_slide_shapes(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

    # Header
    add_text(slide, "ANSWER 04 · YES, FOR NOW", emu(0.7), emu(0.3), emu(10), emu(0.35),
             size=13, bold=True, color=RED, letter_spacing=400)
    add_text(slide, "Yes — they do different jobs.",
             emu(0.7), emu(0.7), emu(12), emu(0.7),
             size=28, bold=True, color=NAVY)

    # Two-column comparison table
    tbl_top = emu(1.8)
    col_w = emu(5.6)
    col_h = emu(1.8)
    gap = emu(0.3)
    x1 = emu(0.7); x2 = x1 + col_w + gap

    # Figma column
    add_rounded_rect(slide, x1, tbl_top, col_w, col_h, fill=LGT)
    add_text(slide, "FIGMA", x1 + emu(0.3), tbl_top + emu(0.15), col_w - emu(0.6), emu(0.3),
             size=13, bold=True, color=GRAY, letter_spacing=300)
    add_text(slide, "Design tool", x1 + emu(0.3), tbl_top + emu(0.45), col_w - emu(0.6), emu(0.4),
             size=20, bold=True, color=NAVY)
    add_text(slide, "Ideation · detail polish · pixel-level craft",
             x1 + emu(0.3), tbl_top + emu(0.9), col_w - emu(0.6), emu(0.35),
             size=13, color=NAVY)
    add_text(slide, "Shared component libraries · plugins · community",
             x1 + emu(0.3), tbl_top + emu(1.25), col_w - emu(0.6), emu(0.35),
             size=13, color=GRAY)

    # Claude Design column
    add_rounded_rect(slide, x2, tbl_top, col_w, col_h, fill=SOFT_RED)
    add_text(slide, "CLAUDE DESIGN", x2 + emu(0.3), tbl_top + emu(0.15), col_w - emu(0.6), emu(0.3),
             size=13, bold=True, color=RED, letter_spacing=300)
    add_text(slide, "Design production tool", x2 + emu(0.3), tbl_top + emu(0.45), col_w - emu(0.6), emu(0.4),
             size=20, bold=True, color=NAVY)
    add_text(slide, "Fast mockups · prototypes · slide / one-pager output",
             x2 + emu(0.3), tbl_top + emu(0.9), col_w - emu(0.6), emu(0.35),
             size=13, color=NAVY)
    add_text(slide, "AI-generated scaffolding · one-click variations",
             x2 + emu(0.3), tbl_top + emu(1.25), col_w - emu(0.6), emu(0.35),
             size=13, color=GRAY)

    # MCP bridge diagram — row 1 (EXPORT), nodes row, row (IMPORT) labels
    # EXPORT label (above the diagram, red)
    add_text(slide, "EXPORT  →  Claude Design output into Figma for final polish",
             emu(0.7), emu(3.95), emu(11.0), emu(0.3),
             size=12, bold=True, color=RED, align='center')

    diag_top = emu(4.35)
    diag_h = emu(1.6)
    # LEFT node: Claude Code + Claude Design
    n1_x = emu(0.7); n1_w = emu(3.5)
    add_rounded_rect(slide, n1_x, diag_top, n1_w, diag_h, fill=WHITE, line=NAVY, line_pt=2)
    # mini terminal
    t_x = n1_x + emu(0.3); t_y = diag_top + emu(0.25); t_w = n1_w - emu(0.6); t_h = emu(1.0)
    add_rounded_rect(slide, t_x, t_y, t_w, t_h, fill=DARK_TERM)
    for i, c in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
        add_oval(slide, t_x + emu(0.15) + i*emu(0.18), t_y + emu(0.15), emu(0.06), fill=c)
    add_text(slide, "$ claude", t_x + emu(0.12), t_y + emu(0.3), t_w - emu(0.25), emu(0.22),
             size=11, color=TERM_GREEN, font='Menlo')
    add_text(slide, "● designing…", t_x + emu(0.12), t_y + emu(0.55), t_w - emu(0.25), emu(0.22),
             size=10, color=TERM_DIM, font='Menlo')
    add_text(slide, "Claude Code + Claude Design", n1_x, diag_top + diag_h - emu(0.4),
             n1_w, emu(0.35), size=14, bold=True, color=NAVY, align='center')

    # CENTER: MCP bridge pill
    n2_x = emu(5.0); n2_w = emu(2.35); n2_y = diag_top + emu(0.3); n2_h = diag_h - emu(0.6)
    add_rounded_rect(slide, n2_x, n2_y, n2_w, n2_h, fill=RED)
    # plug icon (simplified)
    plug_cx = n2_x + emu(0.45); plug_cy = n2_y + n2_h/2
    add_oval(slide, plug_cx, plug_cy, emu(0.22), fill=WHITE)
    add_rect(slide, plug_cx - emu(0.16), plug_cy - emu(0.04), emu(0.1), emu(0.05), fill=RED)
    add_rect(slide, plug_cx + emu(0.06), plug_cy - emu(0.04), emu(0.1), emu(0.05), fill=RED)
    # MCP label (no letter_spacing — LibreOffice underline artifact)
    add_text(slide, "MCP", plug_cx + emu(0.3), plug_cy - emu(0.3), n2_w - emu(0.8), emu(0.4),
             size=20, bold=True, color=WHITE)
    add_text(slide, "bridge", plug_cx + emu(0.3), plug_cy + emu(0.05), n2_w - emu(0.8), emu(0.3),
             size=12, color=WHITE)

    # RIGHT node: Figma file
    n3_x = emu(7.95); n3_w = emu(3.65)
    add_rounded_rect(slide, n3_x, diag_top, n3_w, diag_h, fill=WHITE, line=NAVY, line_pt=2)
    # Figma logo (simplified 5-shape F at top-center of node)
    fig_cx = n3_x + n3_w/2; fig_cy = diag_top + emu(0.25)
    sz = emu(0.18)
    # Top row: orange circle + purple square
    add_oval(slide, fig_cx - sz, fig_cy + sz, sz*0.85, fill=RGBColor(0xF2, 0x4E, 0x1E))
    add_rect(slide, fig_cx, fig_cy, sz*1.7, sz*1.7, fill=RGBColor(0xA2, 0x59, 0xFF))
    # Middle: blue square + green circle
    add_rect(slide, fig_cx - sz*1.7, fig_cy + sz*1.7, sz*1.7, sz*1.7, fill=RGBColor(0x1A, 0xBC, 0xFE))
    add_oval(slide, fig_cx, fig_cy + sz*2.55, sz*0.85, fill=RGBColor(0x0A, 0xCF, 0x83))
    # Bottom: coral circle
    add_oval(slide, fig_cx - sz, fig_cy + sz*4.2, sz*0.85, fill=RGBColor(0xFF, 0x72, 0x62))
    add_text(slide, "Figma file", n3_x, diag_top + diag_h - emu(0.55),
             n3_w, emu(0.35), size=14, bold=True, color=NAVY, align='center')
    add_text(slide, "ideation · polish", n3_x, diag_top + diag_h - emu(0.25),
             n3_w, emu(0.22), size=11, color=GRAY, align='center')

    # Arrows — simple straight bars from each node to center MCP on the "top track" and back on "bottom track"
    top_arr_y = diag_top + emu(0.35)
    bot_arr_y = diag_top + diag_h - emu(0.35)
    # top bars (left→right, red)
    add_rect(slide, n1_x + n1_w - emu(0.02), top_arr_y - emu(0.025), n2_x - (n1_x + n1_w) + emu(0.04), emu(0.05), fill=RED)
    add_rect(slide, n2_x + n2_w - emu(0.02), top_arr_y - emu(0.025), n3_x - (n2_x + n2_w) + emu(0.04), emu(0.05), fill=RED)
    # bottom bars (right→left, navy)
    add_rect(slide, n1_x + n1_w - emu(0.02), bot_arr_y - emu(0.025), n2_x - (n1_x + n1_w) + emu(0.04), emu(0.05), fill=NAVY)
    add_rect(slide, n2_x + n2_w - emu(0.02), bot_arr_y - emu(0.025), n3_x - (n2_x + n2_w) + emu(0.04), emu(0.05), fill=NAVY)

    # IMPORT label (below diagram, navy)
    add_text(slide, "IMPORT  ←  Legacy Figma project → context for Claude Code",
             emu(0.7), diag_top + diag_h + emu(0.1), emu(11.0), emu(0.3),
             size=12, bold=True, color=NAVY, align='center')

    # Footer note
    add_text(slide, "No matter how tech grows, painters still need pens and canvas.",
             emu(0.7), emu(6.75), emu(12), emu(0.4),
             size=14, color=GRAY, align='center')

def build_conclusion_2(slide):
    """Conclusion 2 — scenario-specific recommendation for the team.

    Draws from claude-design-vs-cli-designer-handoff.md's "針對這場景的建議"
    section: two scenario cards + a relay callout.
    """
    clear_slide_shapes(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

    # Header
    add_text(slide, "CONCLUSION · FOR OUR SCENARIO", emu(0.7), emu(0.35), emu(10), emu(0.4),
             size=14, bold=True, color=RED)
    add_text(slide, "Pick the tool by the job.",
             emu(0.7), emu(0.85), emu(12), emu(0.85),
             size=36, bold=True, color=NAVY)

    # Two scenario cards
    card_top = emu(2.05)
    card_h = emu(3.5)
    card_w = emu(5.65)
    gap = emu(0.3)
    x1 = emu(0.7); x2 = x1 + card_w + gap

    # --- Card 1: Visual validation / cross-team review ---
    add_rounded_rect(slide, x1, card_top, card_w, card_h, fill=LGT)
    # icon badge
    add_rounded_rect(slide, x1 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35), fill=NAVY)
    add_text(slide, "A", x1 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35),
             size=14, bold=True, color=WHITE, align='center', valign='middle')
    add_text(slide, "Pure visual validation · cross-team review",
             x1 + emu(1.05), card_top + emu(0.3), card_w - emu(1.4), emu(0.4),
             size=13, bold=True, color=GRAY)
    # Headline
    add_text(slide, "Claude Design\nis the main stage.",
             x1 + emu(0.35), card_top + emu(0.95), card_w - emu(0.7), emu(1.3),
             size=26, bold=True, color=NAVY)
    # Body bullets
    bullets_a = [
        "Org-internal share URL for stakeholder review",
        "Inline comments and live canvas iteration",
        "One-click export to PDF / PPTX / Canva",
    ]
    for i, b in enumerate(bullets_a):
        by = card_top + emu(2.3 + i*0.35)
        add_rounded_rect(slide, x1 + emu(0.4), by + emu(0.08), emu(0.12), emu(0.12), fill=NAVY)
        add_text(slide, b, x1 + emu(0.65), by, card_w - emu(0.95), emu(0.3),
                 size=13, color=NAVY)

    # --- Card 2: Ship to RD / automate / version control ---
    add_rounded_rect(slide, x2, card_top, card_w, card_h, fill=SOFT_RED)
    add_rounded_rect(slide, x2 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35), fill=RED)
    add_text(slide, "B", x2 + emu(0.35), card_top + emu(0.35), emu(0.55), emu(0.35),
             size=14, bold=True, color=WHITE, align='center', valign='middle')
    add_text(slide, "Ship to RD · automate requirements · version design system",
             x2 + emu(1.05), card_top + emu(0.3), card_w - emu(1.4), emu(0.4),
             size=13, bold=True, color=RED)
    # Headline
    add_text(slide, "Claude Code CLI leads.\nClaude Design scouts ahead.",
             x2 + emu(0.35), card_top + emu(0.95), card_w - emu(0.7), emu(1.3),
             size=22, bold=True, color=NAVY)
    # Body bullets
    bullets_b = [
        "Main: Claude Code CLI — automation pipeline, one fewer breakpoint",
        "Claude Design up front — designer ideation + stakeholder review",
        "Then handoff bundle → Claude Code → RD",
    ]
    for i, b in enumerate(bullets_b):
        by = card_top + emu(2.3 + i*0.35)
        add_rounded_rect(slide, x2 + emu(0.4), by + emu(0.08), emu(0.12), emu(0.12), fill=RED)
        add_text(slide, b, x2 + emu(0.65), by, card_w - emu(0.95), emu(0.3),
                 size=12, color=NAVY)

    # Relay callout strip at bottom
    relay_y = emu(5.85)
    relay_h = emu(1.1)
    add_rounded_rect(slide, emu(0.7), relay_y, emu(11.9), relay_h, fill=NAVY)
    # Left: emphasis statement
    add_text(slide, "It's not either / or — it's a relay.",
             emu(1.0), relay_y + emu(0.18), emu(5.8), emu(0.45),
             size=20, bold=True, color=WHITE)
    add_text(slide, "Two tools. One chain. Different jobs.",
             emu(1.0), relay_y + emu(0.65), emu(5.8), emu(0.35),
             size=13, color=RGBColor(0xB8, 0xC3, 0xC9))
    # Right: chain visualization
    chain_y = relay_y + emu(0.38)
    # Design pill
    pill_w = emu(1.3); pill_h = emu(0.38); pill_gap = emu(0.35)
    cx1 = emu(7.3)
    cx2 = cx1 + pill_w + pill_gap
    cx3 = cx2 + pill_w + pill_gap
    for label, x, color in [
        ("Design", cx1, RED),
        ("Code", cx2, RGBColor(0x30, 0x5A, 0x72)),
        ("RD", cx3, RGBColor(0x78, 0x80, 0x86)),
    ]:
        add_rounded_rect(slide, x, chain_y, pill_w, pill_h, fill=color)
        add_text(slide, label, x, chain_y, pill_w, pill_h,
                 size=13, bold=True, color=WHITE, align='center', valign='middle')
    # arrows between pills
    for ax in [cx1 + pill_w, cx2 + pill_w]:
        add_text(slide, "→", ax, chain_y, pill_gap, pill_h,
                 size=18, color=WHITE, align='center', valign='middle')

def build_references(slide):
    clear_slide_shapes(slide)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    add_text(slide, "APPENDIX", emu(0.7), emu(0.5), emu(6), emu(0.4),
             size=13, bold=True, color=RED, letter_spacing=400)
    add_text(slide, "References", emu(0.7), emu(0.9), emu(10), emu(0.8),
             size=36, bold=True, color=NAVY)

    refs = [
        ("1", "Anthropic.", "Introducing Claude Design by Anthropic Labs.",
         "anthropic.com/news/claude-design-anthropic-labs"),
        ("2", "Claude (@claudeai).", "Launch post on X, April 2026.",
         "x.com/claudeai/status/2045156267690213649"),
        ("3", "Ryan Mather (@Flomerboy).", "Tips for getting the best results out of Claude Design.",
         "x.com/Flomerboy/status/2045162321589252458"),
        ("4", "Anthropic.", "Claude Design (product).",
         "claude.ai/design"),
        ("5", "Figma.", "Figma Skills: Code to design workflows.",
         "figma.com/community/skills"),
        ("6", "Rasmus Andersson (@rsms).", "Quote on design vs. design-production tools.",
         "x.com/rsms/status/2045239193971179851"),
    ]
    col1_x = emu(0.7); col2_x = emu(6.8)
    col_w = emu(5.7)
    row_h = emu(1.4)
    for i, (num, who, title, url) in enumerate(refs):
        col = 0 if i < 3 else 1
        row = i if i < 3 else i - 3
        x = col1_x if col == 0 else col2_x
        y = emu(2.0) + row * row_h
        # Number badge
        add_rounded_rect(slide, x, y, emu(0.4), emu(0.4), fill=RED)
        add_text(slide, num, x, y, emu(0.4), emu(0.4),
                 size=14, bold=True, color=WHITE, align='center', valign='middle')
        # Content
        tx = x + emu(0.55); tw = col_w - emu(0.55)
        add_text(slide, who, tx, y, tw, emu(0.3),
                 size=12, bold=True, color=NAVY)
        add_text(slide, title, tx, y + emu(0.3), tw, emu(0.55),
                 size=12, color=NAVY)
        add_text(slide, url, tx, y + emu(0.9), tw, emu(0.3),
                 size=10, color=GRAY, font='Menlo')

# ---------- Main ----------
def main():
    prs = Presentation(SRC)

    # The v1 has 22 slides (indices 0..21). We'll add 11 new slides at the end
    # (indices 22..32), then drop old #3 (idx 2) and old #18 (idx 17), then
    # reorder.
    print(f"Loaded v1: {len(prs.slides)} slides")

    # --- Create new slides ---
    # Use 'Black Blank' (idx 7) as the base for all new slides (clean canvas,
    # we'll add our own shapes).
    new_slides = {}

    s = add_blank_slide(prs, 7); build_toc(s);               new_slides['toc'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_subtitle(s);          new_slides['subtitle'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_q_prompt(s, "Q1", "What's the difference between\nClaude Code and Claude Design?");  new_slides['q1'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_a1_callout(s);        new_slides['a1'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_q_prompt(s, "Q2", "How can Claude Design help us?");           new_slides['q2'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_a2_text_callout(s);   new_slides['a2'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_q_prompt(s, "Q3", "Does Claude Design improve\nthe quality of design delivery?"); new_slides['q3'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_a3_callout(s);        new_slides['a3'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_q_prompt(s, "Q4", "Do we still need Figma?");  new_slides['q4'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_a4_callout(s);        new_slides['a4'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_conclusion_2(s);      new_slides['concl2'] = len(prs.slides) - 1
    s = add_blank_slide(prs, 7); build_references(s);        new_slides['refs'] = len(prs.slides) - 1

    print(f"After insert: {len(prs.slides)} slides. New indices: {new_slides}")

    # --- Build reorder mapping ---
    # Old indices (0-based):
    #  0=cover, 1=What's Claude Design, 2=What's the Difference (drop),
    #  3=Pain Points, 4=In short, 5=What CD Solved, 6=How Compares,
    #  7=Hands-on URL, 8=Design System, 9=Hands-on subtitle,
    #  10=DS import, 11=DS import result, 12=Start to design, 13=Sketch canva,
    #  14=Before starting, 15=Before starting dup, 16=Outcome,
    #  17=Further Questions (drop), 18=Figma official skills,
    #  19=What's Next, 20=Conclusion, 21=Ryan Mather quote
    t = new_slides
    order = [
        0,                       # 1  Title cover
        1,                       # 2  Briefing (What's Claude Design)
        t['toc'],                # 3  TOC (new)
        t['subtitle'],           # 4  Subtitle "Let's start..."
        t['q1'],                 # 5  Q1 prompt
        t['a1'],                 # 6  A1 callout (GitHub/Git analogy)
        6,                       # 7  How Claude Design Compares (old #7, idx 6)
        t['q2'],                 # 8  Q2 prompt
        t['a2'],                 # 9  A2 text callout
        3,                       # 10 Pain Points (old #4, idx 3)
        4,                       # 11 In short (old #5, idx 4)
        5,                       # 12 What Claude Design Solved (old #6, idx 5)
        t['q3'],                 # 13 Q3 prompt
        t['a3'],                 # 14 A3 callout (camera/car/claude)
        t['q4'],                 # 15 Q4 prompt
        t['a4'],                 # 16 A4 callout (Figma MCP flow)
        18,                      # 17 Figma official skills (old #19, idx 18)
        7,                       # 18 Hands-on URL (old #8, idx 7)
        8,                       # 19 Design System (old #9, idx 8)
        9,                       # 20 Hands-on subtitle (old #10, idx 9)
        10,                      # 21 Design system import (old #11, idx 10)
        11,                      # 22 DS import result (old #12, idx 11)
        12,                      # 23 Start to design (old #13, idx 12)
        13,                      # 24 Sketch canva (old #14, idx 13)
        14,                      # 25 Before starting (old #15, idx 14)
        15,                      # 26 Before starting dup (old #16, idx 15)
        16,                      # 27 Outcome (old #17, idx 16)
        20,                      # 28 Conclusion (old #21, idx 20) — generic benefit/can't-solve
        t['concl2'],             # 29 Conclusion 2 (new) — scenario-specific recommendation
        19,                      # 30 What's Next (old #20, idx 19)
        21,                      # 31 Ryan Mather quote (old #22, idx 21)
        t['refs'],               # 32 References (new)
    ]

    # Drop old indices 2 (What's the Difference) and 17 (Further Questions):
    # we achieve this simply by not including them in `order`. The slides
    # remain in the package but won't be referenced by sldIdLst, so they won't
    # show. Clean removal would require also deleting the slide parts, but
    # sldIdLst omission is enough to hide them from viewers.
    # (A future cleanup can fully delete via prs.part.drop_rel + XML surgery.)

    print(f"Reorder length: {len(order)}")
    reorder_slides(prs, order)

    prs.save(DST)
    print(f"Saved: {DST}")
    print(f"Final slide count (visible): {len(order)}")

if __name__ == "__main__":
    main()
