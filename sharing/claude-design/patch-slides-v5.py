"""Round 5 — post-企業帳號驗證 polish + Slide 32 Update 內容填入.

Changes:
  #5  類比從 "Claude Design : Claude Code = GitHub : Git"
      改成 "Claude Design : Claude Code ≈ GitHub Desktop : GitHub"
      訊息從 "same underlying work" 改成「GUI 是子集、能力更局限」
      左下對照標籤 "GITHUB · GUI" → "GITHUB DESKTOP · GUI"
  #15 "Design production tool" → "Design delivery tool"
      （避免 RD 混合聽眾把 "production" 誤解成 production code）
  #32 填入 anti-ai-slop 摘要 + Claude Design vs huashu-design 對照表
      （保留現有的 Title "Update" + 副標 + GitHub URL）

Run in place:
  python3 patch-slides-v5.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "claude-design-sharing.pptx")

# ---------- Text replacements (Slide 5 + 15) ----------
REPLACEMENTS = [
    # Slide 5: GitHub Desktop : GitHub analogy
    (4, "Claude Design : Claude Code  =  GitHub : Git",
        "Claude Design : Claude Code  ≈  GitHub Desktop : GitHub"),
    (4, "A GUI layer that directs the same underlying work.",
        "A simplified GUI covering the daily 80% — the rest stays in the full tool."),
    (4, "GITHUB   ·   GUI", "GITHUB DESKTOP   ·   GUI"),
    (4, "GIT   ·   CLI", "GITHUB   ·   FULL"),
    # Slide 15: delivery tool (avoid RD misreading "production")
    (14, "Design production tool", "Design delivery tool"),
]

TEMPLATE_BLUE = RGBColor(0x1F, 0x4E, 0x79)   # eyebrow blue (matches deck)
TEMPLATE_GRAY = RGBColor(0x55, 0x66, 0x77)
TEMPLATE_DARK = RGBColor(0x15, 0x32, 0x41)
ACCENT_LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)
BORDER_GRAY   = RGBColor(0xD0, 0xD7, 0xDE)


def apply_run_level(slide, old, new):
    replaced = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if old not in full:
                continue
            hit = False
            for r in para.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    replaced += 1
                    hit = True
                    break
            if hit:
                continue
            first = para.runs[0]
            first.text = full.replace(old, new, 1)
            for r in para.runs[1:]:
                r.text = ""
            replaced += 1
    return replaced


def fill_slide32(slide):
    """Add summary paragraph + comparison table, preserving existing title/subtitle/URL."""
    # Layout constants. 16:9 slide is 13.333 × 7.5 inch (12192000 × 6858000 EMU).
    slide_width_in  = 13.333
    slide_height_in = 7.5

    # Summary paragraph — placed under the existing subtitle area.
    summary_left   = Inches(0.6)
    summary_top    = Inches(1.75)
    summary_width  = Inches(slide_width_in - 1.2)
    summary_height = Inches(1.15)

    tb = slide.shapes.add_textbox(summary_left, summary_top, summary_width, summary_height)
    tf = tb.text_frame
    tf.word_wrap = True

    # Line 1: what is AI slop
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r = p1.add_run()
    r.text = "AI slop = AI 預設產出的「視覺最大公約數」"
    r.font.name = "Noto Sans TC"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = TEMPLATE_DARK

    r2 = p1.add_run()
    r2.text = " — 紫漸層、emoji icon、圓角卡片、AI 手畫 SVG。"
    r2.font.name = "Noto Sans TC"
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEMPLATE_DARK

    # Line 2
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(4)
    r = p2.add_run()
    r.text = "問題不是醜，是"
    r.font.name = "Noto Sans TC"
    r.font.size = Pt(13)
    r.font.color.rgb = TEMPLATE_GRAY
    r2 = p2.add_run()
    r2.text = "不攜帶任何品牌資訊"
    r2.font.name = "Noto Sans TC"
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.color.rgb = TEMPLATE_DARK
    r3 = p2.add_run()
    r3.text = "——每個 AI 做出來的頁面都長一樣。"
    r3.font.name = "Noto Sans TC"
    r3.font.size = Pt(13)
    r3.font.color.rgb = TEMPLATE_GRAY

    # Line 3 — thesis
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(6)
    r = p3.add_run()
    r.text = "Claude Design 與 huashu-design 都在反 slop，解法殊途同歸："
    r.font.name = "Noto Sans TC"
    r.font.size = Pt(13)
    r.font.italic = True
    r.font.color.rgb = TEMPLATE_BLUE
    r2 = p3.add_run()
    r2.text = " 前者靠「建 system」，後者靠「抓真實資產」。"
    r2.font.name = "Noto Sans TC"
    r2.font.size = Pt(13)
    r2.font.italic = True
    r2.font.color.rgb = TEMPLATE_BLUE

    # --- Comparison table ---
    rows = 5
    cols = 3
    tbl_left   = Inches(0.6)
    tbl_top    = Inches(3.05)
    tbl_width  = Inches(slide_width_in - 1.2)
    tbl_height = Inches(3.7)

    table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
    tbl = table_shape.table

    # Column widths (sum = 12.133 inches)
    col_widths = [Inches(2.5), Inches(4.8), Inches(4.833)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    headers = ["維度", "Claude Design", "huashu-design"]
    data = [
        ("形態",
         "網頁產品（GUI）· 訂閱 quota",
         "Skill（Claude Code 對話）· API 消耗，可並行跑 agent"),
        ("遇到空白時",
         "引導用戶建一套 design system（選色、字、component 規則）",
         "進入「設計方向顧問模式」——從 20 種設計哲學推 3 個差異化方向"),
        ("有具體品牌時",
         "仍走 system 流程",
         "跳過 system，走 5 步硬流程：搜官方 → 下載 logo / 產品圖 / UI → 寫 brand-spec.md"),
        ("核心賭注",
         "規範優先（build system first）",
         "資產優先（fetch assets first）— logo / 產品圖 > 色值 / 字體"),
    ]

    # Header row
    for c, text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.name = "Noto Sans TC"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Header cell fill
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEMPLATE_BLUE
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)

    # Data rows
    for r_idx, (dim, cd, hua) in enumerate(data, start=1):
        for c_idx, text in enumerate([dim, cd, hua]):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = text
            r.font.name = "Noto Sans TC"
            r.font.size = Pt(11)
            if c_idx == 0:
                r.font.bold = True
                r.font.color.rgb = TEMPLATE_DARK
            else:
                r.font.color.rgb = TEMPLATE_DARK
            # Alternating row tint
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_LIGHT if r_idx % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)

    # Row heights
    tbl.rows[0].height = Inches(0.45)
    for i in range(1, rows):
        tbl.rows[i].height = Inches((tbl_height.inches - 0.45) / (rows - 1))

    # Footer soundbite
    ft_left   = Inches(0.6)
    ft_top    = Inches(6.85)
    ft_width  = Inches(slide_width_in - 1.2)
    ft_height = Inches(0.55)
    ft = slide.shapes.add_textbox(ft_left, ft_top, ft_width, ft_height)
    ftf = ft.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.LEFT
    fr = fp.add_run()
    fr.text = "只要 agent 是「為某個具體品牌」工作，而不是「為 AI 平均值」工作，slop 就不會發生。"
    fr.font.name = "Noto Sans TC"
    fr.font.size = Pt(11)
    fr.font.italic = True
    fr.font.color.rgb = TEMPLATE_GRAY


def main():
    prs = Presentation(SRC)

    # Text replacements on slides 5 + 15
    for idx, old, new in REPLACEMENTS:
        slide = prs.slides[idx]
        n = apply_run_level(slide, old, new)
        status = "OK" if n > 0 else "MISS"
        print(f"  [{status}] slide {idx+1}: {old[:60]!r} → {new[:60]!r}  ({n} hit)")

    # Slide 32 content fill (index 31)
    slide32 = prs.slides[31]
    fill_slide32(slide32)
    print(f"  [OK] slide 32: summary + comparison table + footer added")

    prs.save(SRC)
    print(f"\nSaved: {SRC}")


if __name__ == "__main__":
    main()
