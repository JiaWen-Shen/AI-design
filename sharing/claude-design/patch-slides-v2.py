"""Round 2 patch. Modifies slides 13, 28 of current pptx.

Preserves user's existing edits on all other slides, and on the Card B
limits of slide 28 (which user has already customized).

Slide 13 (A3 callout) — append a "Why?" block explaining Claude Design's
  harness + collaboration weaknesses leading to inability to deliver designs
  consistently across time and teams.

Slide 28 (Conclusion 2) — add a 'team-project 可共編' bullet to Card A,
  and replace the bottom relay section with a pair of scenario lists
  mapping Card A / Card B.

Run:
  cd .../sharing/claude-design/
  python3 patch-slides-v2.py
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from importlib.machinery import SourceFileLoader
bm = SourceFileLoader(
    "buildmod",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "build-claude-design-sharing.py")
).load_module()

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "claude-design-sharing.pptx")
DST = SRC

RED = bm.RED; NAVY = bm.NAVY; GRAY = bm.GRAY
LGT = bm.LGT; WHITE = bm.WHITE; SOFT_RED = bm.SOFT_RED
emu = bm.emu
SLIDE_W = bm.SLIDE_W; SLIDE_H = bm.SLIDE_H

# ---------- Slide 13 (A3) — add a Why? block at the bottom ----------
def add_why_block_to_a3(slide):
    """Shrink existing caption strip height and append a Why? explanation box."""
    # Find and remove the existing bottom caption strip ("Easier to drive ..."
    # text plus its pink pill). We'll redraw a more compact version, and slot
    # the Why block underneath.
    to_remove = []
    for sh in slide.shapes:
        try:
            # caption strip: rounded rect at y ~6.8 with SOFT_RED fill
            if hasattr(sh, 'fill') and sh.fill.type == 1:
                rgb = sh.fill.fore_color.rgb
                if rgb and str(rgb).lower() == 'fbe5e6' and sh.top and sh.top > emu(6.5):
                    to_remove.append(sh)
            # caption text at same y range
            if sh.has_text_frame and sh.top and sh.top > emu(6.5):
                t = sh.text_frame.text
                if 'Easier to drive' in t or 'Wider range' in t:
                    to_remove.append(sh)
        except Exception:
            pass
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)

    # Redraw caption (compact) at y=6.65 in a thinner pill
    cap_y = emu(6.6)
    bm.add_rounded_rect(slide, emu(1.5), cap_y, emu(10.3), emu(0.4), fill=SOFT_RED)
    bm.add_text(slide, "Easier to drive.    ·    Wider range.",
                emu(1.5), cap_y, emu(10.3), emu(0.4),
                size=14, bold=True, color=NAVY, align='center', valign='middle')

    # Why? block — RED flag + 2-line explanation + consequence line
    why_y = emu(7.1)
    why_h = emu(0.35)
    bm.add_text(slide, "WHY?", emu(0.7), why_y, emu(1.0), why_h,
                size=12, bold=True, color=RED, valign='middle')
    bm.add_text(slide,
                "Claude Design's harness (no persistent docs / git / versioned tokens) and collaboration (no skill-share / submodule / decision trace) stay thin — so it cannot consistently deliver designs across time and across teams.",
                emu(1.6), why_y, emu(11.0), why_h,
                size=11, color=NAVY, valign='middle')

# ---------- Slide 28 (Conclusion 2) — add team-edit bullet + scenario lists ----------
def patch_conclusion_2(slide):
    """Add 'team-project collaboration' bullet to Card A, replace bottom
    relay section with Card A / Card B scenario mapping.
    User's Card B edits (4 limits, 'Better harness' subtitle) are preserved.
    """
    # Card A: add a 4th bullet
    # Find the last Card A bullet (the square before "One-click export ...")
    # and its text, then add a new bullet below.
    # Card A lives roughly x=0.7-6.35. Bullets are at x~1.10 (dots) and x~1.35+ (text).
    # Simpler: just draw a new bullet at known position.
    card_top = emu(1.85)
    x1 = emu(0.7); card_w = emu(5.65)
    # Position: same vertical rhythm as existing 3 bullets (at 2.25, 2.60, 2.95).
    # Add 4th at y = card_top + emu(3.3)
    by = card_top + emu(3.3)
    bm.add_rounded_rect(slide, x1 + emu(0.4), by + emu(0.08), emu(0.12), emu(0.12), fill=NAVY)
    bm.add_text(slide, "Team-project workspaces — shared folders, real-time co-edit (Team / Enterprise plan)",
                x1 + emu(0.65), by, card_w - emu(0.95), emu(0.3),
                size=12, color=NAVY)

    # Remove the existing bottom relay strip (NAVY rounded rect + its inner
    # shapes). We identify the NAVY rect by fill + y>=5.8.
    to_remove = []
    for sh in slide.shapes:
        try:
            if sh.top and sh.top >= emu(5.8):
                # include ALL shapes below the cards (relay + chain pills + labels)
                to_remove.append(sh)
        except Exception:
            pass
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)

    # Redraw the bottom as a two-column scenario map.
    # Header strip (thin, just says relay + Design→Code→RD chain)
    bot_top = emu(5.85)
    # Small header bar
    bm.add_text(slide, "IT'S A RELAY — Design → Code → RD",
                emu(0.7), bot_top, emu(11.9), emu(0.35),
                size=12, bold=True, color=RED, align='left')
    # Horizontal divider
    bm.add_line(slide, emu(0.7), bot_top + emu(0.4), emu(12.6), bot_top + emu(0.4),
                rgb=GRAY, width_pt=0.5)

    # Two scenario columns, one under each card above
    col_y = bot_top + emu(0.5)
    col_h = emu(1.0)
    col_w = emu(5.65)
    gap = emu(0.3)
    cx1 = emu(0.7); cx2 = cx1 + col_w + gap

    # --- Column under Card A ---
    bm.add_text(slide, "A · USE CLAUDE DESIGN WHEN",
                cx1, col_y, col_w, emu(0.28),
                size=11, bold=True, color=NAVY)
    scenarios_a = [
        "Concept ideation · stakeholder review decks",
        "Marketing one-pagers · pitch artifacts",
        "Cross-team sign-off on visual direction",
    ]
    for i, s in enumerate(scenarios_a):
        y = col_y + emu(0.3) + i * emu(0.23)
        bm.add_rounded_rect(slide, cx1, y + emu(0.08), emu(0.08), emu(0.08), fill=NAVY)
        bm.add_text(slide, s, cx1 + emu(0.18), y, col_w - emu(0.2), emu(0.23),
                    size=11, color=NAVY)

    # --- Column under Card B ---
    bm.add_text(slide, "B · FALL BACK TO CLAUDE CODE WHEN",
                cx2, col_y, col_w, emu(0.28),
                size=11, bold=True, color=RED)
    scenarios_b = [
        "Production handoff to RD — with a repeatable pipeline",
        "PM requirements flowing from Jira / Teams / Confluence",
        "Versioned design system shared across squads",
    ]
    for i, s in enumerate(scenarios_b):
        y = col_y + emu(0.3) + i * emu(0.23)
        bm.add_rounded_rect(slide, cx2, y + emu(0.08), emu(0.08), emu(0.08), fill=RED)
        bm.add_text(slide, s, cx2 + emu(0.18), y, col_w - emu(0.2), emu(0.23),
                    size=11, color=NAVY)

# ---------- Main ----------
def main():
    prs = Presentation(SRC)
    print(f"Loaded: {len(prs.slides)} slides")

    # Slide 13 is index 12 (A3 callout)
    add_why_block_to_a3(prs.slides[12])
    print("  patched slide 13 (A3 → added Why? block)")

    # Slide 28 is index 27 (Conclusion 2)
    patch_conclusion_2(prs.slides[27])
    print("  patched slide 28 (Conclusion 2 → +team-edit bullet, scenario map at bottom)")

    prs.save(DST)
    print(f"Saved: {DST}")

if __name__ == "__main__":
    main()
