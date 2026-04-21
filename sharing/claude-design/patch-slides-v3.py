"""Round 3 patch.

Slide 13 — restructure the WHY text block for readability. User moved it
  to a portrait area on the left (x=0.70, y=1.97, w=3.73, h=3.03). Convert
  from one run-on paragraph into a clear 3-section structure.

Slide 28 — light polish to user's edited Card B scenarios.
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from importlib.machinery import SourceFileLoader
bm = SourceFileLoader(
    "buildmod",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "build-claude-design-sharing.py")
).load_module()

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "claude-design-sharing.pptx")
DST = SRC

RED = bm.RED; NAVY = bm.NAVY; GRAY = bm.GRAY
LGT = bm.LGT; WHITE = bm.WHITE; SOFT_RED = bm.SOFT_RED
emu = bm.emu

# ---------- Slide 13 — rewrite WHY text block for readability ----------
def refine_why_block(slide):
    # Find the why text box by content signature
    target = None
    for sh in slide.shapes:
        if sh.has_text_frame and 'harness' in sh.text_frame.text and 'collaboration' in sh.text_frame.text:
            target = sh
            break
    if target is None:
        print("  WARN: could not find WHY text box")
        return

    # Preserve existing box geometry
    left, top, w, h = target.left, target.top, target.width, target.height
    # Wipe the body — rebuild paragraphs
    tf = target.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = True
    tf.clear()

    # Helper to set paragraph with runs of varying format
    def para(parts, space_before_pt=0, space_after_pt=0, align_left=True):
        """parts: list of (text, size, bold, color, italic=False)."""
        # Use first paragraph if fresh, else add
        if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align_left:
            p.alignment = PP_ALIGN.LEFT
        if space_before_pt:
            p.space_before = Pt(space_before_pt)
        if space_after_pt:
            p.space_after = Pt(space_after_pt)
        for part in parts:
            text, size, bold, color = part[:4]
            italic = part[4] if len(part) > 4 else False
            r = p.add_run()
            r.text = text
            r.font.name = 'Inter'
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            r.font.underline = False

    # --- Build structured WHY ---
    # Row 1: "Harness stays thin" headline + details line
    para([("Harness stays thin.", 18, True, NAVY)], space_after_pt=2)
    para([("No persistent docs.  No git.  No versioned tokens.", 12, False, GRAY)],
         space_after_pt=14)

    # Row 2: "Collaboration stays shallow" headline + details line
    para([("Collaboration stays shallow.", 18, True, NAVY)], space_after_pt=2)
    para([("No skill share.  No submodule.  No decision trace.", 12, False, GRAY)],
         space_after_pt=14)

    # Row 3: consequence
    para([("→  ", 14, True, RED),
          ("So it can't consistently deliver designs", 14, True, NAVY)])
    para([("across time, or across teams.", 14, True, NAVY)])

# ---------- Slide 28 — polish Card B scenarios (user-edited) ----------
def refine_conclusion2_scenarios(slide):
    """User wrote:
      'Cross function collaboration'
      'Importing live requirements, design systems, policies…etc'
    Polish into consistent, tighter phrasing.
    """
    # Find Card B scenario text box (bottom-right, y > 6.0, text contains 'Cross')
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            if 'Cross' in t and 'collaboration' in t.lower():
                # Replace with polished 2-line version
                tf = sh.text_frame
                tf.clear()
                def para(text, size=11, bold=False, color=NAVY, space_after_pt=4):
                    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.space_after = Pt(space_after_pt)
                    r = p.add_run()
                    r.text = text
                    r.font.name = 'Inter'
                    r.font.size = Pt(size)
                    r.font.bold = bold
                    r.font.color.rgb = color
                    r.font.underline = False
                para("Cross-function traceable collaboration — PM ↔ design ↔ RD")
                para("Live ingest of requirements, design system, policy updates")
                print(f"  refined Card B scenarios (was: {t[:80]!r})")
                return
    print("  WARN: couldn't find Card B scenarios box")

# ---------- Main ----------
def main():
    prs = Presentation(SRC)
    print(f"Loaded: {len(prs.slides)} slides")

    refine_why_block(prs.slides[12])
    print("  refined slide 13 WHY block (structured 3-section layout)")

    refine_conclusion2_scenarios(prs.slides[27])

    prs.save(DST)
    print(f"Saved: {DST}")

if __name__ == "__main__":
    main()
